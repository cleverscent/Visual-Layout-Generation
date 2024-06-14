import json
import pickle
import os
import numpy as np 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import concurrent.futures
import torch
import clip
from transformers import AutoTokenizer, AutoModel, AutoModelForSeq2SeqLM
from PIL import Image
from tqdm import tqdm

# JSON 파일 불러오기
def load_json_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:  # 인코딩 명시적으로 지정
        data = json.load(file)
    return data

# JSON 파일 저장
def save_json(data, file_path):
    with open(file_path, 'w', encoding='utf-8') as file:
        json.dump(data, file, indent=4)

########################### PPT parsing #############################
def shape_type_to_string(shape_type):
    """
    MSO_SHAPE_TYPE 값을 문자열로 변환합니다.
    """
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "auto_shape"
    elif shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text_box"
    elif shape_type == MSO_SHAPE_TYPE.CALLOUT:
        return "callout"
    elif shape_type == MSO_SHAPE_TYPE.CANVAS:
        return "canvas"
    elif shape_type == MSO_SHAPE_TYPE.CHART:
        return "chart"
    elif shape_type == MSO_SHAPE_TYPE.COMMENT:
        return "comment"
    elif shape_type == MSO_SHAPE_TYPE.DIAGRAM:
        return "diagram"
    elif shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        return "embedded_ole_object"
    elif shape_type == MSO_SHAPE_TYPE.FORM_CONTROL:
        return "form_control"
    elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
        return "freeform"
    elif shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    elif shape_type == MSO_SHAPE_TYPE.IGX_GRAPHIC:
        return "smartart_graphic"
    elif shape_type == MSO_SHAPE_TYPE.INK:
        return "ink"
    elif shape_type == MSO_SHAPE_TYPE.INK_COMMENT:
        return "ink_comment"
    elif shape_type == MSO_SHAPE_TYPE.LINE:
        return "line"
    elif shape_type == MSO_SHAPE_TYPE.LINKED_OLE_OBJECT:
        return "linked_ole_object"
    elif shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
        return "linked_picture"
    elif shape_type == MSO_SHAPE_TYPE.MEDIA:
        return "media"
    elif shape_type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
        return "ole_control_object"
    elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "placeholder"
    elif shape_type == MSO_SHAPE_TYPE.SCRIPT_ANCHOR:
        return "script_anchor"
    elif shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"
    elif shape_type == MSO_SHAPE_TYPE.TEXT_EFFECT:
        return "text_effect"
    elif shape_type == MSO_SHAPE_TYPE.WEB_VIDEO:
        return "web_video"
    elif shape_type == MSO_SHAPE_TYPE.MIXED:
        return "mixed"
    else:
        return "other"
    
def process_shape_info(shape, slide_number, z_index, ppt_name, image_base_path):
    rotation = shape.rotation
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    
    # shape_type_to_string 함수를 사용하여 shape의 type을 문자열로 변환
    shape_type_str = shape_type_to_string(shape.shape_type)
    
    # Check if the shape has text content
    text_content = shape.text_frame.text if shape.has_text_frame else None
    is_text = text_content is not None and text_content.strip() != ""  # Check if text content is not an empty string
    
    # 이미지 파일 이름 생성
    image_file_name = f"Slide{slide_number}_Shape{z_index+1}.png"
    image_path = os.path.join(image_base_path, ppt_name, image_file_name)
    
    # 이미지의 실제 너비와 높이를 가져오기
    try:
        with Image.open(image_path) as img:
            img_width, img_height = img.size  # 이미지의 실제 너비와 높이
    except FileNotFoundError:
        print('############### image not found! ################')
        print(image_path)
        img_width, img_height = 0, 0  # 이미지 파일이 없을 경우

    return {
        "type": shape_type_str,
        "rotation": rotation,
        "left": left / 914400 * 96,  # EMU to pixel
        "top": top / 914400 * 96,  # EMU to pixel
        "width": width / 914400 * 96,  # EMU to pixel
        "height": height / 914400 * 96,  # EMU to pixel
        "img_width": img_width,  # 이미지의 실제 너비 추가
        "img_height": img_height,  # 이미지의 실제 높이 추가
        "is_text": is_text,
        "text_content": None if not is_text else text_content,  # Set text_content to None if is_text is False
        "slide_number": slide_number,
        "z_index": z_index,
        "image_file_name": image_file_name,
    }


def process_slide(slide, page_number, ppt_name, image_base_path):
    slide_data = {"contents": []}
    for z_index, shape in enumerate(slide.shapes):
        content_info = process_shape_info(shape, slide_number=page_number, z_index=z_index, ppt_name=ppt_name,image_base_path=image_base_path)
        slide_data["contents"].append(content_info)
    return slide_data

def process_pptx(input_pptx_folder, output_json, image_base_path):
    all_ppt_data = {"presentations": []}
    all_shape_types = set()  # 모든 shape 타입들을 저장할 집합
    
    for ppt_file in os.listdir(input_pptx_folder):
        if ppt_file.endswith(".pptx"):
            presentation_path = os.path.join(input_pptx_folder, ppt_file)
            presentation = Presentation(presentation_path)
            ppt_name = os.path.splitext(ppt_file)[0]

            result_data = {
                "ppt_name": ppt_name,
                "slides": [],
                "presentation_size": {
                    "width": presentation.slide_width.pt*4/3,
                    "height": presentation.slide_height.pt*4/3
                }
            }

            for i, slide in enumerate(presentation.slides, start=1):
                slide_data = process_slide(slide, page_number=i, ppt_name=ppt_name, image_base_path=image_base_path)
                result_data["slides"].append(slide_data)
                
                # 각 shape의 타입을 all_shape_types 집합에 추가
                for content in slide_data["contents"]:
                    all_shape_types.add(content["type"])

            all_ppt_data["presentations"].append(result_data)

    with open(output_json, "w", encoding="utf-8") as json_file:
        json.dump(all_ppt_data, json_file, indent=2, ensure_ascii=False)
        
    print("All parsed shape types:", all_shape_types)
        

################################# 요소 20개 이상 slide 제거 ###########################################

def remove_dense_slides(input_json, output_json):
    presentations = load_json_file(input_json)
    
    for presentation in presentations['presentations']:
        presentation['slides'] = [slide for slide in presentation['slides'] if len(slide['contents']) < 41]
        
    save_json(presentations, output_json)

############################## CLIP 모델을 사용한 이미지 처리 ###########################################

# 이미지 파일 경로와 ppt 이름 추출 (경로 구성 방식 업데이트)
def extract_image_paths_and_ppt_names(json_data, base_path):
    image_info = []
    for presentation in json_data['presentations']:
        ppt_name = presentation['ppt_name']
        ppt_base_path = os.path.join(base_path, ppt_name)
        for slide in presentation['slides']:
            for content in slide['contents']:
                image_path = os.path.join(ppt_base_path, content['image_file_name'])
                image_info.append((image_path, ppt_name))
    return image_info


        
def process_images_with_clip(input_json, base_path, output_json):
    json_data = load_json_file(input_json)
    # CLIP 모델 로드
    device = "cuda" if torch.cuda.is_available() else "cpu"
    model, preprocess = clip.load("ViT-B/32", device=device)
    
    image_info = extract_image_paths_and_ppt_names(json_data, base_path)
    
    # 이미지 처리 함수 수정
    def process_image(image_path):
        try:
            image = Image.open(image_path).convert("RGBA").resize((224, 224))
            image = preprocess(image).unsqueeze(0).to(device)
            with torch.no_grad():
                image_features = model.encode_image(image)
            return image_features.cpu().numpy().tolist()
        except FileNotFoundError:
            print(f"File not found: {image_path}")
            return None

    # 병렬 처리로 이미지 CLIP 모델에 입력
    results = {}
    with concurrent.futures.ThreadPoolExecutor() as executor:
        future_to_image = {executor.submit(process_image, info[0]): info for info in image_info}
        
        for future in tqdm(concurrent.futures.as_completed(future_to_image), total=len(image_info), desc="Processing images"):
            image_path, ppt_name = future_to_image[future]
            image_features = future.result()
            if image_features:
                if ppt_name not in results:
                    results[ppt_name] = {}
                results[ppt_name][os.path.basename(image_path)] = image_features
    

    # 결과 저장
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=4)
        

############################## Text encoder을 사용한 text 처리 ###########################################

def load_json_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        data = json.load(file)
    return data

def extract_text_content_and_slide_names(json_data):
    extracted_texts = []

    for presentation in json_data["presentations"]:
        for slide in presentation["slides"]:
            for content in slide["contents"]:
                if not content.get("is_text", False) or not content.get("text_content"):
                    # If is_text is False or text_content is not available, treat as null
                    content_id = f"{presentation['ppt_name']}/{content['image_file_name']}"
                    extracted_texts.append({
                        "presentation_name": presentation['ppt_name'],
                        "image_file_name": content['image_file_name'],
                        "text_content": None,
                        "content_id": content_id,
                    })
                else:
                    content_id = f"{presentation['ppt_name']}/{content['image_file_name']}"
                    extracted_texts.append({
                        "presentation_name": presentation['ppt_name'],
                        "image_file_name": content['image_file_name'],
                        "text_content": content["text_content"],
                        "content_id": content_id,
                    })

    return extracted_texts

def encode_text_content(text_content, tokenizer, model):
    if text_content is None:
        # Return null value for text_content
        return [0] * 1024  # Assuming 768 dimensions for null value
    else:
        tokens = tokenizer(text_content, return_tensors="pt", padding=True, truncation=True)
        with torch.no_grad():
            output = model(**tokens)
        return output.last_hidden_state.mean(dim=1).squeeze().tolist()

def process_text_features(input_json, output_json):
    # Load the JSON data
    json_data = load_json_file(input_json)

    # Load the T5 tokenizer and model
    tokenizer = AutoTokenizer.from_pretrained("t5--base-multilingual")
    model = AutoModelForSeq2SeqLM.from_pretrained("t5-xxl-base-multilingual")

    # Process and encode text features
    text_features = []

    text_info = extract_text_content_and_slide_names(json_data)

    for info in tqdm(text_info, desc="Processing text features"):
        text_feature = encode_text_content(info["text_content"], tokenizer, model)
        text_features.append({
            "presentation_name": info["presentation_name"],
            "image_file_name": info["image_file_name"],
            "text_feature": text_feature
        })

    # Save the text features to a new JSON file with the extracted text features
    output_data = {}
    for item in text_features:
        presentation_name = item["presentation_name"]
        image_file_name = item["image_file_name"]

        if presentation_name not in output_data:
            output_data[presentation_name] = {}

        output_data[presentation_name][image_file_name] = [item["text_feature"]]

    with open(output_json, "w", encoding="utf-8") as output_file:
        json.dump(output_data, output_file, indent=4)


################################################## classical clustering method ##################################################
from sklearn.cluster import KMeans 
from sklearn.cluster import AgglomerativeClustering 
from sklearn.mixture import GaussianMixture 
from sklearn.cluster import DBSCAN 
from sklearn.decomposition import PCA 

def cluster(data_json_path, model_output_path, num_pca_components, num_clusters, eps, min_samples, clustering_method):
    with open(data_json_path, 'r', encoding='utf-8') as f:
        json_data = json.load(f)
    
    # Get data from json file
    data = []
    for ppt_name, slide_data in json_data.items():
        for image_file_name, image_feature_list in slide_data.items():
            for feature in image_feature_list:
                data.append(feature)

    # Convert data to NumPy array
    data = np.array(data)

    # Apply PCA 
    pca = PCA(n_components=num_pca_components)
    
    # Choose clustering method
    if clustering_method == 'kmeans':
        clustering_model = KMeans(n_clusters=num_clusters, random_state=42)
    elif clustering_method == 'hierarchical': # You can try 'average' linkage too
        clustering_model = AgglomerativeClustering(n_clusters=num_clusters, metric='euclidean', linkage='complete')  
    elif clustering_method == 'gmm':
        clustering_model = GaussianMixture(n_components=num_clusters, random_state=42)
    elif clustering_method == 'dbscan':
        clustering_model = DBSCAN(eps=eps, min_samples=min_samples)
    else:
        raise ValueError("Invalid clustering method. Choose from 'kmeans', 'hierarchical', 'gmm', or 'dbscan'.")


    # PCA와 KMeans 모델을 dictionary에 저장
    model_dict = {'pca': pca, 'clustering_model': clustering_model}

    # 클러스터링 모델을 .pkl 파일로 저장
    with open(model_output_path, 'wb') as file:
        pickle.dump(model_dict, file)

################################################## ppt 슬라이드 이미지로 캡쳐하기 ##################################################
import os
import fitz
from pptx import Presentation
from comtypes import client

def pptx_to_pdf_comtypes(pptx_path, pdf_path):
    powerpoint = client.CreateObject('PowerPoint.Application')
    presentation = powerpoint.Presentations.Open(pptx_path)
    presentation.SaveAs(pdf_path, FileFormat=32)  # 32 = PDF 형식
    presentation.Close()
    powerpoint.Quit()

def pptx_to_pdf(input_folder, output_folder):
    # 모든 pptx 파일에 대해 반복
    for filename in os.listdir(input_folder):
        if filename.endswith(".pptx"):
            pptx_path = os.path.join(input_folder + "\\" + filename)
            ppt_name = os.path.splitext(filename)[0]  # 확장자를 제외한 파일 이름

            # pdf를 저장할 폴더 생성
            pdf_folder = os.path.join(output_folder + "\\" + ppt_name)
            os.makedirs(pdf_folder, exist_ok=True)

            # pptx를 pdf로 변환
            pdf_path = os.path.join(pdf_folder + "\\" + f"{ppt_name}.pdf")
            pptx_to_pdf_comtypes(pptx_path, pdf_path)

            # pdf에서 이미지로 저장
            pdf_doc = fitz.open(pdf_path)
            for i, page in enumerate(pdf_doc):
                img = page.get_pixmap()
                img.save(os.path.join(pdf_folder + "\\" + f"{i}.png"))
            pdf_doc.close()
            
            # pdf 파일 삭제
            os.remove(pdf_path)


###########################################################################################################

# 설정 파일 불러오기 및 처리 실행
def run_from_config(config_path):
    config = load_json_file(config_path)
    if config["process_pptx"]:
        process_pptx(config["input_pptx_folder"], config["output_json"], config["image_base_path"])

    if config["remove_dense_slides"]:
        remove_dense_slides(config["input_json_for_removal"], config["output_json_after_removal"])

    if config["process_images_with_clip"]:
        process_images_with_clip(config["input_json_for_clip"], config["base_path_for_images"], config["ouput_json_after_clip"])
    
    if config["process_texts"]:
        process_text_features(config["input_json_for_text_encoder"], config["ouput_json_after_text_encoder"])
        
    if config["cluster_with_image"]:
        cluster(config["image_json_path"], config["clusterd_image_model_path"], config["num_pca_components"],
                config["num_clusters"], config["eps"], config["min_samples"], config["clustering_method"])

    if config["cluster_with_text"]:
        cluster(config["text_json_path"], config["clusterd_text_model_path"], config["num_pca_components"],
                config["num_clusters"], config["eps"], config["min_samples"], config["clustering_method"])
    
    if config["Capture_all_slides_in_a_PPT_as_images"]:
        config["input_folder"] = config["input_folder"].replace("/", "\\")
        config["output_folder"] = config["output_folder"].replace("/", "\\")
        pptx_to_pdf(config["input_folder"], config["output_folder"])

###########################################################################################################
if __name__ == "__main__":
    config_path = "./data_loaders/preprocessing_config.json"  # 설정 파일 경로
    run_from_config(config_path)
